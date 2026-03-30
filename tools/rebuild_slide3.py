#!/usr/bin/env python3
"""Replace slides 3 and 4 with a single combined 'The Solution' slide."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

INPUT = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations", "MBSE_Pipeline_Executive_Brief.pptx")


def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox


def add_text(tf, text, size=11, bold=False, color=(0xCC, 0xCC, 0xCC), space_after=4, alignment=PP_ALIGN.LEFT, first=False):
    if first and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    p.space_after = Pt(space_after)
    p.alignment = alignment
    return p


def add_rounded_box(slide, left, top, width, height, fill_rgb, border_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    if border_rgb:
        shape.line.color.rgb = RGBColor(*border_rgb)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_stage_box(slide, left, top, width, height, title, desc, fill_rgb, text_rgb=(0xFF, 0xFF, 0xFF)):
    """Add a pipeline stage box with title and description."""
    box = add_rounded_box(slide, left, top, width, height, fill_rgb)
    txBox = add_textbox(slide, left + Inches(0.1), top + Inches(0.08), width - Inches(0.2), height - Inches(0.1))
    tf = txBox.text_frame
    add_text(tf, title, size=9, bold=True, color=text_rgb, space_after=1, first=True)
    add_text(tf, desc, size=7, bold=False, color=(0xBB, 0xBB, 0xBB), space_after=0)
    return box


def add_arrow(slide, left, top, width, height):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x69, 0xB0)
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height):
    """Add a down-pointing arrow."""
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x69, 0xB0)
    shape.line.fill.background()
    return shape


def main():
    prs = Presentation(INPUT)
    print(f"Before: {len(prs.slides)} slides")

    # Delete slides 3 and 4 (indices 2 and 3)
    # Delete index 3 first (so index 2 doesn't shift)
    for idx in [3, 2]:
        rId = prs.slides._sldIdLst[idx].get(qn('r:id'))
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[idx]

    print(f"After delete: {len(prs.slides)} slides")

    # Add new combined slide
    slide_layout = prs.slides[0].slide_layout
    slide = prs.slides.add_slide(slide_layout)

    # Move to position 3 (index 2)
    new_sldId = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(new_sldId)
    prs.slides._sldIdLst.insert(2, new_sldId)

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    # ============================================
    # TITLE + SUBTITLE
    # ============================================
    title_box = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5))
    add_text(title_box.text_frame, "The Solution: Two Pipelines, One Thread", size=28, bold=True,
             color=(0xFF, 0xFF, 0xFF), space_after=0, first=True)

    sub_box = add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4))
    add_text(sub_box.text_frame,
             "A single automated thread connects the Cameo model to the deployed product \u2014 with quality gates enforcing traceability at every stage.",
             size=11, color=(0xAA, 0xAA, 0xAA), space_after=0, first=True)

    # ============================================
    # TOP LANE: Model Pipeline
    # ============================================
    lane_y = Inches(1.8)
    lane_h = Inches(1.8)

    # Lane background
    add_rounded_box(slide, Inches(0.5), lane_y - Inches(0.1), Inches(12.3), lane_h + Inches(0.2),
                    fill_rgb=(0x14, 0x25, 0x38), border_rgb=(0x2A, 0x3A, 0x4A))

    # Lane label
    label_box = add_textbox(slide, Inches(0.7), lane_y, Inches(2.5), Inches(0.35))
    add_text(label_box.text_frame, "MODEL PIPELINE", size=9, bold=True,
             color=(0x00, 0x99, 0xD6), space_after=0, first=True)

    owner_box = add_textbox(slide, Inches(0.7), lane_y + Inches(0.28), Inches(3), Inches(0.25))
    add_text(owner_box.text_frame, "cameo-model-pipeline  \u2022  Owned by Systems Engineers",
             size=8, color=(0x88, 0x88, 0x88), space_after=0, first=True)

    # Stage boxes - top lane
    stage_w = Inches(2.3)
    stage_h = Inches(0.85)
    stage_y = lane_y + Inches(0.7)
    gap = Inches(0.15)

    stages_top = [
        ("Cameo Model", "Author requirements\n& ICD in the model", (0x00, 0x30, 0x57)),
        ("\u2192  Export & Validate", "Schema validation\nof JSON exports", (0x00, 0x4A, 0x87)),
        ("\u2192  Package & Version", "Semantic versioning\n& artifact assembly", (0x00, 0x69, 0xB0)),
        ("\u2192  Publish Artifact", "Versioned ZIP to\nGitHub Releases", (0x00, 0x82, 0xCC)),
    ]

    x = Inches(0.7)
    for title, desc, fill_c in stages_top:
        add_stage_box(slide, x, stage_y, stage_w, stage_h, title, desc, fill_c)
        x += stage_w + gap

    # ============================================
    # DOWN ARROW (versioned artifact)
    # ============================================
    arrow_x = Inches(5.8)
    arrow_y = lane_y + lane_h + Inches(0.15)
    add_down_arrow(slide, arrow_x, arrow_y, Inches(0.5), Inches(0.55))

    artifact_label = add_textbox(slide, arrow_x + Inches(0.6), arrow_y + Inches(0.05), Inches(2.5), Inches(0.4))
    add_text(artifact_label.text_frame, "Versioned Artifact", size=9, bold=True,
             color=(0x00, 0x99, 0xD6), space_after=0, first=True)
    add_text(artifact_label.text_frame, "requirements.json + proto files", size=7,
             color=(0x88, 0x88, 0x88), space_after=0)

    # ============================================
    # BOTTOM LANE: Product Pipeline
    # ============================================
    lane2_y = arrow_y + Inches(0.75)
    lane2_h = Inches(1.8)

    # Lane background
    add_rounded_box(slide, Inches(0.5), lane2_y - Inches(0.1), Inches(12.3), lane2_h + Inches(0.2),
                    fill_rgb=(0x14, 0x25, 0x38), border_rgb=(0x2A, 0x3A, 0x4A))

    # Lane label
    label_box2 = add_textbox(slide, Inches(0.7), lane2_y, Inches(2.5), Inches(0.35))
    add_text(label_box2.text_frame, "PRODUCT PIPELINE", size=9, bold=True,
             color=(0x00, 0x99, 0xD6), space_after=0, first=True)

    owner_box2 = add_textbox(slide, Inches(0.7), lane2_y + Inches(0.28), Inches(3.5), Inches(0.25))
    add_text(owner_box2.text_frame, "product-pipeline  \u2022  Owned by Developers + Systems Engineers",
             size=8, color=(0x88, 0x88, 0x88), space_after=0, first=True)

    # Stage boxes - bottom lane
    stage2_y = lane2_y + Inches(0.7)

    stages_bottom = [
        ("Fetch Model", "Download versioned\nartifact from release", (0x00, 0x30, 0x57)),
        ("\u2192  Quality Gates", "Coverage, drift &\norphan detection", (0xB4, 0x53, 0x09)),
        ("\u2192  BDD Verification", "Automated Gherkin\nlog-analysis tests", (0x00, 0x69, 0xB0)),
        ("\u2192  Deploy + Report", "Helm deploy with\nevidence bundle", (0x00, 0x82, 0xCC)),
    ]

    x = Inches(0.7)
    for title, desc, fill_c in stages_bottom:
        add_stage_box(slide, x, stage2_y, stage_w, stage_h, title, desc, fill_c)
        x += stage_w + gap

    # ============================================
    # KEY PRINCIPLE
    # ============================================
    principle_y = lane2_y + lane2_h + Inches(0.3)
    principle_box = add_rounded_box(slide, Inches(0.5), principle_y, Inches(12.3), Inches(0.5),
                                    fill_rgb=(0x0A, 0x14, 0x20), border_rgb=(0x00, 0x69, 0xB0))

    p_text = add_textbox(slide, Inches(0.8), principle_y + Inches(0.08), Inches(11.8), Inches(0.4))
    add_text(p_text.text_frame,
             "Every requirement is automatically traced from model to test to deployment. No manual assembly. No gaps. No drift.",
             size=11, bold=True, color=(0xFF, 0xFF, 0xFF), space_after=0, alignment=PP_ALIGN.CENTER, first=True)

    # ============================================
    # TOOL STACK (right margin, small)
    # ============================================
    tools_box = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2.3), Inches(0.5))
    add_text(tools_box.text_frame, "Cameo \u2022 Python \u2022 Behave \u2022 Gherkin \u2022 Helm \u2022 GitHub Actions",
             size=7, color=(0x66, 0x66, 0x66), space_after=0, alignment=PP_ALIGN.RIGHT, first=True)

    # ============================================
    # PAGE NUMBER
    # ============================================
    pg_box = add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3))
    add_text(pg_box.text_frame, "3 / 9", size=9, color=(0x88, 0x88, 0x88),
             space_after=0, alignment=PP_ALIGN.RIGHT, first=True)

    # Update page numbers on remaining slides
    print(f"Final: {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Fix old page numbers
                        old_numbers = {
                            "3 / 10": f"{i+1} / {len(prs.slides)}",
                            "4 / 10": f"{i+1} / {len(prs.slides)}",
                            "5 / 10": f"{i+1} / {len(prs.slides)}",
                            "6 / 10": f"{i+1} / {len(prs.slides)}",
                            "7 / 10": f"{i+1} / {len(prs.slides)}",
                            "8 / 10": f"{i+1} / {len(prs.slides)}",
                            "9 / 10": f"{i+1} / {len(prs.slides)}",
                            "10 / 10": f"{i+1} / {len(prs.slides)}",
                            "1 / 10": f"1 / {len(prs.slides)}",
                            "2 / 10": f"2 / {len(prs.slides)}",
                        }
                        for old, new in old_numbers.items():
                            if run.text.strip() == old:
                                run.text = new

    prs.save(INPUT)
    print("Saved!")

    # Verify
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and len(t) > 5:
                        texts.append(t)
                        break
                if texts:
                    break
        print(f"  Slide {i+1}: {texts[0][:60] if texts else '(empty)'}")


if __name__ == "__main__":
    main()
