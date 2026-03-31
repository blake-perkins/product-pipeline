#!/usr/bin/env python3
"""Build the executive briefing deck from scratch as a new file.

Instead of modifying the existing .pptx (which corrupts the zip when
deleting slides), this creates a brand new file with all 11 slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations", "MBSE_Pipeline_Executive_Brief.pptx")

# Brand colors
NAVY = (0x0D, 0x1B, 0x2A)
BLUE_NG = (0x00, 0x69, 0xB0)
BLUE_LIGHT = (0x00, 0x99, 0xD6)
BLUE_4 = (0x00, 0x82, 0xCC)
NAVY_1 = (0x00, 0x30, 0x57)
NAVY_2 = (0x00, 0x4A, 0x87)
WHITE = (0xFF, 0xFF, 0xFF)
BLACK = (0x1A, 0x1A, 0x1A)
GRAY = (0xAA, 0xAA, 0xAA)
GRAY_LIGHT = (0xBB, 0xBB, 0xBB)
GRAY_DIM = (0x88, 0x88, 0x88)
GRAY_DARK = (0x55, 0x55, 0x55)
AMBER = (0xB4, 0x53, 0x09)
# Dark slide surfaces
SURFACE_DARK = (0x14, 0x25, 0x38)
BORDER_DARK = (0x2A, 0x3A, 0x4A)
# Light slide surfaces
SURFACE_LIGHT = (0xF4, 0xF7, 0xFA)
BORDER_LIGHT = (0xD0, 0xE8, 0xF8)
# Aliases for dark slides (backward compat)
SURFACE = SURFACE_DARK
BORDER = BORDER_DARK


def rgb(color):
    return RGBColor(*color)


def add_tb(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def text(tf, content, size=11, bold=False, color=GRAY, after=4, align=PP_ALIGN.LEFT, first=False):
    if first and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    p.space_after = Pt(after)
    p.alignment = align
    return p


def add_box(slide, left, top, width, height, fill, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if border:
        shape.line.color.rgb = rgb(border)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def dark_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(NAVY)


def white_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(WHITE)


def page_num(slide, num, total):
    tb = add_tb(slide, Inches(12), Inches(7), Inches(1), Inches(0.3))
    text(tb.text_frame, f"{num} / {total}", size=9, color=GRAY_DIM, after=0, align=PP_ALIGN.RIGHT, first=True)


def stage_box(slide, left, top, w, h, title, desc, fill):
    add_box(slide, left, top, w, h, fill)
    tb = add_tb(slide, left + Inches(0.1), top + Inches(0.08), w - Inches(0.2), h - Inches(0.1))
    text(tb.text_frame, title, size=9, bold=True, color=WHITE, after=1, first=True)
    text(tb.text_frame, desc, size=7, color=GRAY_LIGHT, after=0)


TOTAL = 11


def slide_01_title(prs):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    dark_bg(slide)
    tb = add_tb(slide, Inches(1), Inches(2), Inches(11), Inches(1))
    text(tb.text_frame, "Model-Based Systems Engineering Pipeline", size=36, bold=True, color=WHITE, after=8, align=PP_ALIGN.CENTER, first=True)
    text(tb.text_frame, "From Requirements to Deployed Product \u2014 Automated, Traceable, Auditable", size=16, color=GRAY, after=16, align=PP_ALIGN.CENTER)
    text(tb.text_frame, "Executive Briefing", size=14, color=BLUE_LIGHT, after=0, align=PP_ALIGN.CENTER)
    page_num(slide, 1, TOTAL)


def slide_02_problem(prs):
    """One Pipeline. Model to Deployment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    # Title
    tb = add_tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6))
    text(tb.text_frame, "One Pipeline. Model to Deployment.", size=32, bold=True, color=NAVY, after=0, first=True)

    # ============================================
    # LEFT SIDE: Three gap cards (60% width)
    # ============================================
    gaps = [
        ("No Shared Language",
         "SEs work in the model, developers work in code. Intent gets lost at every handoff \u2014 misalignment surfaces at integration."),
        ("No Reliable Handoff",
         "Artifacts shared over email have no version control. When the model changes, nothing alerts the team. Tests go stale silently."),
        ("No Connected Toolchain",
         "Requirements, tests, and deployment live in separate tools. Audit evidence is rebuilt from scratch every review cycle."),
    ]

    y = Inches(1.5)
    for gap_title, gap_desc in gaps:
        add_box(slide, Inches(0.5), y, Inches(7.3), Inches(1.35), SURFACE_LIGHT, BORDER_LIGHT)
        tb = add_tb(slide, Inches(0.8), y + Inches(0.15), Inches(6.8), Inches(1.1))
        text(tb.text_frame, gap_title, size=15, bold=True, color=BLUE_NG, after=4, first=True)
        text(tb.text_frame, gap_desc, size=10, color=GRAY_DARK, after=0)
        y += Inches(1.55)

    # ============================================
    # RIGHT SIDE: Disconnected silos diagram
    # ============================================
    diag_x = Inches(8.5)
    diag_w = Inches(4.0)

    # "Today" label
    tb = add_tb(slide, diag_x, Inches(1.5), diag_w, Inches(0.35))
    text(tb.text_frame, "Today", size=14, bold=True, color=GRAY_DIM, after=0, align=PP_ALIGN.CENTER, first=True)

    # Three silo boxes — stacked vertically, not connected
    silo_w = Inches(2.8)
    silo_h = Inches(0.65)
    silo_x = diag_x + Inches(0.6)
    silos = [
        ("Cameo Model", NAVY_1),
        ("Test Suite", NAVY_2),
        ("Deployment", GRAY_DIM),
    ]

    silo_y = Inches(2.1)
    silo_positions = []
    for label, color in silos:
        add_box(slide, silo_x, silo_y, silo_w, silo_h, color)
        tb = add_tb(slide, silo_x + Inches(0.1), silo_y + Inches(0.12), silo_w - Inches(0.2), silo_h - Inches(0.2))
        text(tb.text_frame, label, size=11, bold=True, color=WHITE, after=0, align=PP_ALIGN.CENTER, first=True)
        silo_positions.append(silo_y)
        silo_y += Inches(1.05)

    # Dashed gap indicators between silos (X marks)
    RED = (0xCC, 0x33, 0x33)
    for i in range(len(silo_positions) - 1):
        gap_y = silo_positions[i] + silo_h + Inches(0.08)
        # X mark shape
        tb = add_tb(slide, silo_x + Inches(1.1), gap_y, Inches(0.6), Inches(0.3))
        text(tb.text_frame, "\u2718", size=16, bold=True, color=RED, after=0, align=PP_ALIGN.CENTER, first=True)

    # "No thread" label under the diagram
    tb = add_tb(slide, diag_x, silo_positions[-1] + silo_h + Inches(0.3), diag_w, Inches(0.5))
    text(tb.text_frame, "No automated thread", size=10, bold=True, color=RED, after=2, align=PP_ALIGN.CENTER, first=True)
    text(tb.text_frame, "Manual handoffs at every stage", size=9, color=GRAY_DIM, after=0, align=PP_ALIGN.CENTER)

    # ============================================
    # Bottom takeaway
    # ============================================
    add_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5), SURFACE_LIGHT, BLUE_NG)
    tb = add_tb(slide, Inches(0.8), Inches(6.36), Inches(11.8), Inches(0.4))
    text(tb.text_frame, "This program moves too fast for manual traceability. The pipeline keeps pace.",
         size=12, bold=True, color=NAVY, after=0, align=PP_ALIGN.CENTER, first=True)

    page_num(slide, 2, TOTAL)


def slide_03_solution(prs):
    """The Solution: Gherkin as the Common Language."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    # Title + subtitle
    tb = add_tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5))
    text(tb.text_frame, "The Solution: Gherkin as the Common Language", size=28, bold=True, color=NAVY, after=0, first=True)

    tb = add_tb(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.35))
    text(tb.text_frame, "Systems Engineers and developers communicate through executable specifications \u2014 Gherkin scenarios that trace directly to Cameo requirements and run as automated tests.",
         size=11, color=GRAY_DARK, after=0, first=True)

    # ============================================
    # THE BRIDGE — Three columns showing the connection
    # ============================================

    # Left column: Cameo Model (SEs)
    add_box(slide, Inches(0.5), Inches(1.7), Inches(3.5), Inches(2.8), SURFACE_LIGHT, BORDER_LIGHT)
    tb = add_tb(slide, Inches(0.7), Inches(1.8), Inches(3.1), Inches(0.3))
    text(tb.text_frame, "Cameo Model", size=14, bold=True, color=NAVY, after=0, first=True)
    tb = add_tb(slide, Inches(0.7), Inches(2.15), Inches(3.1), Inches(0.25))
    text(tb.text_frame, "Owned by Systems Engineers", size=9, color=BLUE_LIGHT, after=0, first=True)

    se_items = [
        "Requirements with verification criteria",
        "ICD definitions and message types",
        "Criteria text drives test intent",
        "Model changes trigger drift detection",
    ]
    tb = add_tb(slide, Inches(0.7), Inches(2.55), Inches(3.1), Inches(1.8))
    for i, item in enumerate(se_items):
        text(tb.text_frame, "\u2022  " + item, size=9, color=GRAY_DARK, after=5, first=(i == 0))

    # Center column: Gherkin (shared)
    add_box(slide, Inches(4.4), Inches(1.5), Inches(4.5), Inches(3.2), BLUE_NG)
    tb = add_tb(slide, Inches(4.6), Inches(1.6), Inches(4.1), Inches(0.3))
    text(tb.text_frame, "Gherkin Specifications", size=16, bold=True, color=WHITE, after=0, align=PP_ALIGN.CENTER, first=True)
    tb = add_tb(slide, Inches(4.6), Inches(1.95), Inches(4.1), Inches(0.25))
    text(tb.text_frame, "Co-authored by SEs + Developers", size=9, color=GRAY_LIGHT, after=0, align=PP_ALIGN.CENTER, first=True)

    # Why Gherkin — value props
    tb = add_tb(slide, Inches(4.7), Inches(2.3), Inches(3.9), Inches(0.9))
    props = [
        "Plain English that both sides can read and validate",
        "Machine-executable \u2014 runs as automated tests on every build",
        "Tagged to requirements \u2014 creates the traceability link",
    ]
    for i, prop in enumerate(props):
        text(tb.text_frame, "\u2022  " + prop, size=8, color=WHITE, after=3, first=(i == 0))

    # Compact Gherkin example
    tb = add_tb(slide, Inches(4.7), Inches(3.2), Inches(3.9), Inches(1.3))
    gherkin_lines = [
        ("@REQ:SYS-REQ-001", True),
        ("Feature: Basic ICD Communications", True),
        ("  @VC:SYS-REQ-001-VC-01 @VER:Test", True),
        ("  Scenario: Valid ICD response", False),
        ("    Given the simulation logs are loaded", False),
        ("    Then the logs should contain", False),
        ('         \"IcdResponse: status=OK\"', False),
    ]
    for i, (line, bold_line) in enumerate(gherkin_lines):
        text(tb.text_frame, line, size=7, bold=bold_line,
             color=WHITE if bold_line else GRAY_LIGHT, after=1, first=(i == 0))

    # Right column: System Verification (Devs)
    add_box(slide, Inches(9.3), Inches(1.7), Inches(3.5), Inches(2.8), SURFACE_LIGHT, BORDER_LIGHT)
    tb = add_tb(slide, Inches(9.5), Inches(1.8), Inches(3.1), Inches(0.3))
    text(tb.text_frame, "System Verification", size=14, bold=True, color=NAVY, after=0, first=True)
    tb = add_tb(slide, Inches(9.5), Inches(2.15), Inches(3.1), Inches(0.25))
    text(tb.text_frame, "Implemented by Developers", size=9, color=BLUE_LIGHT, after=0, first=True)

    dev_items = [
        "Step definitions implement the HOW",
        "Log analysis against the deployed product",
        "E2E verification on every build",
        "Results feed the pipeline report",
    ]
    tb = add_tb(slide, Inches(9.5), Inches(2.55), Inches(3.1), Inches(1.8))
    for i, item in enumerate(dev_items):
        text(tb.text_frame, "\u2022  " + item, size=9, color=GRAY_DARK, after=5, first=(i == 0))

    # Arrows: Cameo -> Gherkin -> Tests
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.05), Inches(2.9), Inches(0.3), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(BLUE_NG)
    shape.line.fill.background()

    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.95), Inches(2.9), Inches(0.3), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(BLUE_NG)
    shape.line.fill.background()

    # ============================================
    # BOTTOM: Two pipeline lanes (compact)
    # ============================================
    tb = add_tb(slide, Inches(0.8), Inches(4.9), Inches(5), Inches(0.3))
    text(tb.text_frame, "How It Flows Through the Pipelines", size=14, bold=True, color=NAVY, after=0, first=True)

    # Model Pipeline (compact)
    add_box(slide, Inches(0.5), Inches(5.3), Inches(6.0), Inches(0.5), SURFACE_LIGHT, BORDER_LIGHT)
    tb = add_tb(slide, Inches(0.6), Inches(5.33), Inches(5.8), Inches(0.4))
    p = text(tb.text_frame, "MODEL PIPELINE", size=8, bold=True, color=BLUE_LIGHT, after=0, first=True)
    r = p.add_run()
    r.text = "    Cameo  \u2192  Export  \u2192  Package  \u2192  Publish Artifact"
    r.font.size = Pt(8)
    r.font.bold = False
    r.font.color.rgb = rgb(GRAY_DARK)

    # Arrow down
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.0), Inches(5.85), Inches(0.3), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(BLUE_NG)
    shape.line.fill.background()

    # Product Pipeline (compact)
    add_box(slide, Inches(0.5), Inches(6.15), Inches(6.0), Inches(0.5), SURFACE_LIGHT, BORDER_LIGHT)
    tb = add_tb(slide, Inches(0.6), Inches(6.18), Inches(5.8), Inches(0.4))
    p = text(tb.text_frame, "PRODUCT PIPELINE", size=8, bold=True, color=BLUE_LIGHT, after=0, first=True)
    r = p.add_run()
    r.text = "    Fetch  \u2192  Quality Gates  \u2192  BDD Tests  \u2192  Deploy + Report"
    r.font.size = Pt(8)
    r.font.bold = False
    r.font.color.rgb = rgb(GRAY_DARK)

    # Right side: key points
    key_points = [
        ("Readable by both sides", "SEs validate intent, developers implement steps."),
        ("Machine-traceable", "Tags (@REQ, @VC) create the link to requirements."),
        ("Pipeline-enforced", "Quality gates ensure every VC has a scenario."),
    ]
    y = Inches(5.15)
    for title, desc in key_points:
        tb = add_tb(slide, Inches(7.0), y, Inches(5.5), Inches(0.45))
        p = text(tb.text_frame, title, size=10, bold=True, color=BLUE_NG, after=0, first=True)
        r = p.add_run()
        r.text = "  \u2014  " + desc
        r.font.size = Pt(9)
        r.font.bold = False
        r.font.color.rgb = rgb(GRAY_DARK)
        y += Inches(0.5)

    # Bottom callout
    add_box(slide, Inches(7.0), Inches(6.7), Inches(5.8), Inches(0.4), (0x0A, 0x14, 0x20), BLUE_NG)
    tb = add_tb(slide, Inches(7.2), Inches(6.74), Inches(5.4), Inches(0.3))
    text(tb.text_frame, "One language. Model to deployment. Both sides own the spec.",
         size=10, bold=True, color=WHITE, after=0, align=PP_ALIGN.CENTER, first=True)

    page_num(slide, 3, TOTAL)


def slide_04_releases(prs):
    """Release Planning & Program Progress."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    # Title
    tb = add_tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5))
    text(tb.text_frame, "Release Planning & Program Progress", size=28, bold=True, color=NAVY, after=0, first=True)

    tb = add_tb(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.35))
    text(tb.text_frame, "Requirements are delivered incrementally across releases. The pipeline tracks what\u2019s in scope, what\u2019s deferred, and what\u2019s shipped.",
         size=11, color=GRAY_DARK, after=0, first=True)

    # Release timeline visualization
    releases = [
        ("1.0.0", "Shipped", "Core ICD &\nHealth Monitoring", BLUE_4, True),
        ("1.2.0", "Current", "Degradation &\nConfiguration Mgmt", BLUE_NG, False),
        ("2.0.0", "Planned", "Resilience, Self-Test\n& Data Logging", GRAY_DIM, False),
    ]

    x = Inches(0.5)
    rw = Inches(3.8)
    for ver, status, desc, color, shipped in releases:
        # Release card
        add_box(slide, x, Inches(1.8), rw, Inches(2.0), SURFACE_LIGHT, BORDER_LIGHT)

        # Version header
        add_box(slide, x, Inches(1.8), rw, Inches(0.45), color)
        tb = add_tb(slide, x + Inches(0.1), Inches(1.83), rw - Inches(0.2), Inches(0.35))
        text(tb.text_frame, f"Release {ver}", size=11, bold=True, color=WHITE, after=0, first=True)

        # Status badge
        tb = add_tb(slide, x + Inches(0.1), Inches(2.35), rw - Inches(0.2), Inches(0.3))
        badge_color = BLUE_LIGHT if shipped else (BLUE_NG if status == "Current" else GRAY_DIM)
        text(tb.text_frame, status.upper(), size=8, bold=True, color=badge_color, after=0, first=True)

        # Description
        tb = add_tb(slide, x + Inches(0.1), Inches(2.65), rw - Inches(0.2), Inches(0.9))
        text(tb.text_frame, desc, size=9, color=GRAY_DARK, after=0, first=True)

        x += rw + Inches(0.25)

    # How it works section
    tb = add_tb(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.3))
    text(tb.text_frame, "How It Works", size=16, bold=True, color=NAVY, after=0, first=True)

    how_items = [
        ("release-plan.json", "A single file maps requirements to release versions. Checked into Git alongside the code."),
        ("Cumulative Scope", "Each release includes all prior releases. Release 1.2.0 enforces everything from 1.0.0 + 1.1.0 + 1.2.0."),
        ("Deferred VCs", "Out-of-scope verification criteria appear as \u201cdeferred\u201d \u2014 visible but not blocking the current release."),
    ]
    y = Inches(4.7)
    for title, desc in how_items:
        tb = add_tb(slide, Inches(0.8), y, Inches(5.5), Inches(0.6))
        text(tb.text_frame, title, size=11, bold=True, color=BLUE_NG, after=2, first=True)
        text(tb.text_frame, desc, size=9, color=GRAY_DARK, after=0)
        y += Inches(0.65)

    # Program progress section (right side)
    tb = add_tb(slide, Inches(7.0), Inches(4.2), Inches(5), Inches(0.3))
    text(tb.text_frame, "Program Visibility", size=16, bold=True, color=NAVY, after=0, first=True)

    viz_items = [
        ("Readiness Ring", "At-a-glance percentage showing how many VCs are passing for the selected release."),
        ("Release Filter", "A global dropdown scopes the entire dashboard to any release \u2014 past, current, or future."),
        ("Release Progress Tab", "Full roadmap with per-release progress bars, target dates, and clickable VC details."),
    ]
    y = Inches(4.7)
    for title, desc in viz_items:
        tb = add_tb(slide, Inches(7.0), y, Inches(5.5), Inches(0.6))
        text(tb.text_frame, title, size=11, bold=True, color=BLUE_NG, after=2, first=True)
        text(tb.text_frame, desc, size=9, color=GRAY_DARK, after=0)
        y += Inches(0.65)

    # Bottom callout
    add_box(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.45), SURFACE_LIGHT, BLUE_NG)
    tb = add_tb(slide, Inches(0.8), Inches(6.66), Inches(11.8), Inches(0.35))
    text(tb.text_frame, "One dashboard shows the full program \u2014 what shipped, what\u2019s in progress, and what\u2019s planned. Updated on every build.",
         size=11, bold=True, color=NAVY, after=0, align=PP_ALIGN.CENTER, first=True)

    page_num(slide, 4, TOTAL)


def slide_05_testing(prs):
    """Layered Testing Approach."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    # Title
    tb = add_tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5))
    text(tb.text_frame, "Where This Fits: Layered Testing", size=28, bold=True, color=NAVY, after=0, first=True)

    tb = add_tb(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.35))
    text(tb.text_frame, "The MBSE pipeline automates E2E testing \u2014 deploying the product, running simulations, and verifying requirements through log analysis.",
         size=11, color=GRAY_DARK, after=0, first=True)

    # Test pyramid — 5 layers matching the standard testing pyramid
    # Each layer: (center_x_offset, y, width, height, label, color, highlighted)
    # Pyramid narrows from bottom to top, centered around x=3.0
    cx = 3.0  # center of pyramid
    layers = [
        (Inches(0.5), Inches(5.4), Inches(5.5), Inches(0.75), "Static Analysis", GRAY_DIM, False),
        (Inches(0.9), Inches(4.5), Inches(4.7), Inches(0.75), "Unit Tests", GRAY_DIM, False),
        (Inches(1.3), Inches(3.6), Inches(3.9), Inches(0.75), "Integration Tests", GRAY_DIM, False),
        (Inches(1.7), Inches(2.7), Inches(3.1), Inches(0.75), "E2E Tests", BLUE_NG, True),
        (Inches(2.1), Inches(1.8), Inches(2.3), Inches(0.75), "Manual &\nExploratory", GRAY_DIM, False),
    ]

    for lx, ly, lw, lh, label, color, highlight in layers:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, ly, lw, lh)
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        if highlight:
            shape.line.color.rgb = rgb(BLUE_LIGHT)
            shape.line.width = Pt(2.5)
        else:
            shape.line.fill.background()

        tb = add_tb(slide, lx + Inches(0.1), ly + Inches(0.1), lw - Inches(0.2), lh - Inches(0.1))
        text(tb.text_frame, label, size=11, bold=True, color=WHITE, after=0, align=PP_ALIGN.CENTER, first=True)

    # Cost arrow (left side)
    tb = add_tb(slide, Inches(0.05), Inches(1.7), Inches(0.4), Inches(4.5))
    text(tb.text_frame, "Cost \u2191", size=9, bold=True, color=GRAY_DIM, after=0, align=PP_ALIGN.CENTER, first=True)

    # Speed arrow (bottom right of pyramid)
    tb = add_tb(slide, Inches(5.5), Inches(6.2), Inches(0.8), Inches(0.3))
    text(tb.text_frame, "Speed \u2192", size=9, bold=True, color=GRAY_DIM, after=0, first=True)

    # "MBSE Pipeline" label pointing to E2E layer
    tb = add_tb(slide, Inches(5.1), Inches(2.8), Inches(1.8), Inches(0.3))
    text(tb.text_frame, "\u25c0 MBSE Pipeline", size=11, bold=True, color=BLUE_NG, after=0, first=True)

    # Right side — what the pipeline does
    tb = add_tb(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(0.3))
    text(tb.text_frame, "What the Pipeline Provides", size=16, bold=True, color=NAVY, after=0, first=True)

    pipeline_items = [
        ("Requirement Traceability", "Every Cameo requirement maps to a Gherkin scenario. Gate A enforces 100% coverage."),
        ("Automated E2E Verification", "BDD scenarios run against simulation logs after product deployment. Results captured automatically."),
        ("Drift Detection", "Gate B detects when the model changes and flags tests for review \u2014 no silent invalidation."),
        ("Evidence Generation", "The Deployment Pipeline Report is produced on every build \u2014 a complete audit trail."),
        ("Release-Aware Gating", "Quality gates only enforce VCs in scope for the current release. Future work is tracked but not blocking."),
    ]
    y = Inches(2.0)
    for title, desc in pipeline_items:
        tb = add_tb(slide, Inches(7.5), y, Inches(5.2), Inches(0.7))
        text(tb.text_frame, title, size=11, bold=True, color=BLUE_NG, after=2, first=True)
        text(tb.text_frame, desc, size=9, color=GRAY_DARK, after=0)
        y += Inches(0.78)

    # Bottom note
    tb = add_tb(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.3))
    text(tb.text_frame, "Static analysis, unit, and integration tests are owned by developers. The MBSE pipeline automates the E2E layer with full requirement traceability.",
         size=9, color=GRAY_DIM, after=0, align=PP_ALIGN.CENTER, first=True)

    page_num(slide, 5, TOTAL)


def slide_06_gates(prs):
    """Three Quality Gates."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    tb = add_tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5))
    text(tb.text_frame, "Three Quality Gates", size=28, bold=True, color=NAVY, after=0, first=True)

    tb = add_tb(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.3))
    text(tb.text_frame, "Automated checks that run on every commit \u2014 the pipeline will not proceed until all gates pass.",
         size=11, color=GRAY_DARK, after=0, first=True)

    gates = [
        ("Gate A: Coverage", BLUE_NG, [
            "Ensures every verification criteria has at least one Gherkin scenario",
            "Generates stub feature files for uncovered VCs",
            "Prevents requirements from reaching deployment without tests",
        ]),
        ("Gate B: Drift Detection", BLUE_4, [
            "Detects when requirement criteria text changes in the model",
            "Compares SHA-256 hash of criteria against baseline",
            "Injects @REVIEW_REQUIRED tag \u2014 forces developer to review & acknowledge",
        ]),
        ("Gate C: Orphan Detection", NAVY_1, [
            "Finds test scenarios referencing non-existent requirements",
            "Catches orphaned @REQ: and @VC: tags",
            "Ensures test suite stays in sync with the model",
        ]),
    ]

    x = Inches(0.5)
    for title, color, items in gates:
        # Gate card
        add_box(slide, x, Inches(1.7), Inches(3.9), Inches(4.5), SURFACE_LIGHT, BORDER_LIGHT)

        # Title bar
        add_box(slide, x, Inches(1.7), Inches(3.9), Inches(0.5), color)
        tb = add_tb(slide, x + Inches(0.15), Inches(1.75), Inches(3.6), Inches(0.4))
        text(tb.text_frame, title, size=14, bold=True, color=WHITE, after=0, first=True)

        # Items
        tb = add_tb(slide, x + Inches(0.2), Inches(2.4), Inches(3.5), Inches(3.5))
        for i, item in enumerate(items):
            text(tb.text_frame, "\u2022  " + item, size=10, color=GRAY_DARK, after=8, first=(i == 0))

        x += Inches(4.1)

    page_num(slide, 4, TOTAL)


def slide_07_tags(prs):
    """Gherkin Specification & Tag Convention."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    tb = add_tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5))
    text(tb.text_frame, "Gherkin Specification & Tag Convention", size=28, bold=True, color=WHITE, after=0, first=True)

    tb = add_tb(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.3))
    text(tb.text_frame, "Tags create the machine-readable traceability link between requirements and test scenarios.",
         size=11, color=GRAY, after=0, first=True)

    # Gherkin example (left side)
    add_box(slide, Inches(0.5), Inches(1.7), Inches(6.5), Inches(4.5), SURFACE, BORDER)
    tb = add_tb(slide, Inches(0.7), Inches(1.85), Inches(6.1), Inches(4.2))
    lines = [
        ("@REQ:SYS-REQ-001", BLUE_LIGHT, True),
        ("Feature: Basic ICD Communications", WHITE, True),
        ("", WHITE, False),
        ("  @VC:SYS-REQ-001-VC-01 @VER:Test", BLUE_LIGHT, True),
        ("  Scenario: Valid ICD produces correct response", WHITE, False),
        ("    Given the simulation logs are loaded", GRAY_LIGHT, False),
        ("    Then the product logs should contain", GRAY_LIGHT, False),
        ('         "Received IcdRequest: test-001"', GRAY_LIGHT, False),
    ]
    for i, (line, color, bold) in enumerate(lines):
        text(tb.text_frame, line if line else " ", size=9, bold=bold, color=color, after=2, first=(i == 0))

    # Tag reference (right side)
    tb = add_tb(slide, Inches(7.3), Inches(1.7), Inches(5.2), Inches(0.3))
    text(tb.text_frame, "Tag Reference", size=14, bold=True, color=WHITE, after=0, first=True)

    tags = [
        ("@REQ:SYS-REQ-001", "Links to requirement ID. Placed at Feature level."),
        ("@VC:...-VC-01", "Links to specific VC. Placed at Scenario level."),
        ("@VER:Test", "Verification criteria type (INCOSE ADIT classification)."),
        ("@REVIEW_REQUIRED", "Injected by Gate B when criteria text drifts."),
        ("@STUB", "Auto-generated placeholder. Replaced by real scenario."),
        ("@DEFERRED", "Injected by pipeline when VC is scoped to a future release."),
    ]
    y = Inches(2.2)
    for tag, desc in tags:
        tb = add_tb(slide, Inches(7.3), y, Inches(5.2), Inches(0.55))
        text(tb.text_frame, tag, size=10, bold=True, color=BLUE_LIGHT, after=2, first=True)
        text(tb.text_frame, desc, size=8, color=GRAY, after=0)
        y += Inches(0.6)

    # Bottom takeaway
    add_box(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.45), SURFACE, BLUE_NG)
    tb = add_tb(slide, Inches(0.8), Inches(6.06), Inches(11.8), Inches(0.35))
    text(tb.text_frame, "One spec. Readable by engineers. Executable by machines. Traceable to every requirement.",
         size=11, bold=True, color=WHITE, after=0, align=PP_ALIGN.CENTER, first=True)

    page_num(slide, 6, TOTAL)


def slide_08_cicd(prs):
    """CI/CD Pipeline Stages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    tb = add_tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5))
    text(tb.text_frame, "CI/CD Pipeline Stages", size=28, bold=True, color=NAVY, after=0, first=True)

    # Model Pipeline
    tb = add_tb(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.3))
    text(tb.text_frame, "Repo 1: Model Pipeline (on every push)", size=12, bold=True, color=NAVY, after=0, first=True)

    model_stages = ["Validate\nExports", "Schema\nCheck", "Generate\nProtos", "Package\nArtifact", "Tag &\nRelease"]
    x = Inches(0.5)
    for s in model_stages:
        stage_box(slide, x, Inches(1.8), Inches(2.2), Inches(0.8), s, "", BLUE_NG)
        x += Inches(2.35)

    # Product Pipeline
    tb = add_tb(slide, Inches(0.8), Inches(3.2), Inches(5), Inches(0.3))
    text(tb.text_frame, "Repo 2: Product Pipeline (on every push)", size=12, bold=True, color=NAVY, after=0, first=True)

    product_stages = [
        ("Traceability\nCheck", BLUE_NG),
        ("Gate A\nCoverage", NAVY_2),
        ("Gate B\nDrift", NAVY_2),
        ("Gate C\nOrphans", NAVY_2),
        ("BDD Test\nExecution", BLUE_4),
        ("Deploy +\nReport", BLUE_4),
    ]
    x = Inches(0.3)
    for s, c in product_stages:
        stage_box(slide, x, Inches(3.7), Inches(1.9), Inches(0.8), s, "", c)
        x += Inches(2.0)

    # Tech stack
    tb = add_tb(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.3))
    text(tb.text_frame, "Technology Stack:  GitHub Actions \u2022 JSON Schema \u2022 Python \u2022 Behave \u2022 Gherkin \u2022 Helm \u2022 Kubernetes",
         size=9, color=GRAY_DIM, after=0, first=True)

    page_num(slide, 7, TOTAL)


def slide_09_cyber(prs):
    """Cyber: SBOM & Vulnerability Scanning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    # Title
    tb = add_tb(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5))
    text(tb.text_frame, "Cyber: Supply Chain Security", size=28, bold=True, color=NAVY, after=0, first=True)

    tb = add_tb(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.35))
    text(tb.text_frame, "Every container image is scanned for known vulnerabilities before deployment. A full Software Bill of Materials is generated on every build.",
         size=11, color=GRAY_DARK, after=0, first=True)

    # LEFT SIDE: Scanning Toolchain (constrained to 6 inches)
    tb = add_tb(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.3))
    text(tb.text_frame, "Automated Scanning Pipeline", size=14, bold=True, color=NAVY, after=0, first=True)

    # Scanning flow — 4 boxes stacked vertically
    scan_steps = [
        ("Container Image", "Product containers built by CI/CD", GRAY_DIM),
        ("Syft \u2192 SBOM", "Generates CycloneDX Software Bill of Materials", BLUE_NG),
        ("Grype \u2192 Scan", "Scans SBOM against known vulnerability databases", BLUE_4),
        ("Policy Gate", "PASS / WARNING / FAIL based on severity thresholds", NAVY_2),
    ]
    y = Inches(2.0)
    for title, desc, color in scan_steps:
        add_box(slide, Inches(0.5), y, Inches(5.8), Inches(0.55), color)
        tb = add_tb(slide, Inches(0.6), y + Inches(0.05), Inches(5.6), Inches(0.45))
        p = text(tb.text_frame, title, size=10, bold=True, color=WHITE, after=0, first=True)
        r = p.add_run()
        r.text = "  \u2014  " + desc
        r.font.size = Pt(8)
        r.font.bold = False
        r.font.color.rgb = rgb(GRAY_LIGHT)
        y += Inches(0.62)

    # What gets produced
    tb = add_tb(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.3))
    text(tb.text_frame, "What Gets Produced", size=14, bold=True, color=NAVY, after=0, first=True)

    artifacts = [
        ("SBOM (CycloneDX)", "Complete component inventory with versions and sources."),
        ("Vulnerability Report", "CVE-level findings with severity and fix availability."),
        ("Policy Verdict", "Blocks deployment on Critical/High findings."),
        ("Severity Dashboard", "Visual breakdown in the Deployment Pipeline Report."),
    ]
    y = Inches(5.0)
    for title, desc in artifacts:
        tb = add_tb(slide, Inches(0.8), y, Inches(5.5), Inches(0.35))
        p = text(tb.text_frame, title, size=10, bold=True, color=BLUE_NG, after=0, first=True)
        r = p.add_run()
        r.text = "  \u2014  " + desc
        r.font.size = Pt(9)
        r.font.bold = False
        r.font.color.rgb = rgb(GRAY_DARK)
        y += Inches(0.38)

    # RIGHT SIDE: Compliance Value
    add_box(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0), SURFACE_LIGHT, BORDER_LIGHT)

    tb = add_tb(slide, Inches(7.0), Inches(1.65), Inches(5.6), Inches(0.3))
    text(tb.text_frame, "Why This Matters for Compliance", size=14, bold=True, color=NAVY, after=0, first=True)

    compliance = [
        ("Executive Order 14028", "Federal mandate requiring SBOMs for all software sold to the government."),
        ("NIST SP 800-218 (SSDF)", "Requires vulnerability scanning and remediation tracking per build."),
        ("Supply Chain Transparency", "Every component inventoried with name, version, type, and package URL."),
        ("Audit-Ready Evidence", "SBOM + scan results included in evidence bundle. No manual assembly."),
        ("Air-Gapped Compatible", "Grype DB mirrored internally. Syft runs against local images. No internet."),
    ]
    y = Inches(2.15)
    for title, desc in compliance:
        tb = add_tb(slide, Inches(7.0), y, Inches(5.6), Inches(0.7))
        text(tb.text_frame, title, size=11, bold=True, color=BLUE_NG, after=2, first=True)
        text(tb.text_frame, desc, size=9, color=GRAY_DARK, after=0)
        y += Inches(0.7)

    # Bottom callout
    add_box(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), SURFACE_LIGHT, BLUE_NG)
    tb = add_tb(slide, Inches(0.8), Inches(6.75), Inches(11.8), Inches(0.3))
    text(tb.text_frame, "Every build produces a complete SBOM and vulnerability scan \u2014 no manual steps, no gaps, always current.",
         size=10, bold=True, color=NAVY, after=0, align=PP_ALIGN.CENTER, first=True)

    page_num(slide, 8, TOTAL)


def slide_10_traceability(prs):
    """Traceability & Audit Evidence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    white_bg(slide)

    tb = add_tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.5))
    text(tb.text_frame, "Traceability & Audit Evidence", size=28, bold=True, color=NAVY, after=0, first=True)

    tb = add_tb(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.3))
    text(tb.text_frame, "Every pipeline run produces a complete, machine-generated evidence chain \u2014 no manual assembly required.",
         size=11, color=GRAY_DARK, after=0, first=True)

    # Flow: Requirement -> VC -> Gherkin -> Test Result -> Report
    flow = ["Requirement\nin Cameo", "Verification\nCriteria (VC)", "Gherkin\nScenario", "BDD Test\nResult", "Deployment\nPipeline Report"]
    x = Inches(0.5)
    for i, label in enumerate(flow):
        c = [NAVY_1, NAVY_2, BLUE_NG, BLUE_4, BLUE_LIGHT][i]
        stage_box(slide, x, Inches(1.7), Inches(2.2), Inches(0.85), label, "", c)
        x += Inches(2.4)

    # Compliance items
    tb = add_tb(slide, Inches(0.8), Inches(3.0), Inches(5), Inches(0.3))
    text(tb.text_frame, "What This Means for Compliance", size=16, bold=True, color=NAVY, after=0, first=True)

    items = [
        ("Minutes, Not Weeks", "Traceability report is generated on every commit \u2014 always current, always complete."),
        ("Every Change Tracked", "Git history + baseline diffs show exactly when requirements changed and who acknowledged them."),
        ("No Missing Tests", "Gate A ensures every verification criteria has coverage \u2014 gaps are caught before merge."),
        ("No Stale Tests", "Gate B detects criteria drift \u2014 existing tests are flagged for re-review when requirements change."),
        ("No Ghost Tests", "Gate C catches orphaned tests \u2014 scenarios for deleted requirements are flagged for removal."),
        ("Immutable Evidence", "CI artifacts (HTML + JSON reports) are stored per build \u2014 reproducible audit trail."),
    ]

    # Two columns
    for col in range(2):
        y = Inches(3.5)
        for i in range(3):
            idx = col * 3 + i
            title, desc = items[idx]
            x = Inches(0.8) if col == 0 else Inches(6.8)
            tb = add_tb(slide, x, y, Inches(5.5), Inches(0.9))
            text(tb.text_frame, title, size=12, bold=True, color=NAVY, after=2, first=True)
            text(tb.text_frame, desc, size=9, color=GRAY_DARK, after=0)
            y += Inches(1.0)

    page_num(slide, 8, TOTAL)


def slide_11_closing(prs):
    """Why This Matters — closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(slide)

    tb = add_tb(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8))
    text(tb.text_frame, "Why This Matters", size=36, bold=True, color=WHITE, after=0, align=PP_ALIGN.CENTER, first=True)

    cols = [
        ("Reduce Risk", "Every requirement is verified before deployment. No gaps, no surprises. Quality gates catch problems early."),
        ("Accelerate Delivery", "Automated traceability eliminates weeks of manual evidence assembly. CI/CD runs on every commit."),
        ("Ensure Compliance", "Machine-generated audit trail from model to deployment. Immutable, reproducible, always current."),
    ]
    for i, (title, desc) in enumerate(cols):
        x = Inches(0.8 + i * 4.1)
        tb = add_tb(slide, x, Inches(2.8), Inches(3.6), Inches(2))
        text(tb.text_frame, title, size=18, bold=True, color=BLUE_LIGHT, after=8, align=PP_ALIGN.CENTER, first=True)
        text(tb.text_frame, desc, size=11, color=GRAY_LIGHT, after=0, align=PP_ALIGN.CENTER)

    tb = add_tb(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4))
    text(tb.text_frame, "One pipeline.  Model to deployment.  Full traceability.  Every commit.",
         size=14, bold=True, color=WHITE, after=0, align=PP_ALIGN.CENTER, first=True)

    tb = add_tb(slide, Inches(4.5), Inches(5.9), Inches(4), Inches(0.5))
    text(tb.text_frame, "Questions?", size=22, bold=True, color=BLUE_LIGHT, after=0, align=PP_ALIGN.CENTER, first=True)

    page_num(slide, 9, TOTAL)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_releases(prs)
    slide_05_testing(prs)
    slide_06_gates(prs)
    slide_07_tags(prs)
    slide_08_cicd(prs)
    slide_09_cyber(prs)
    slide_10_traceability(prs)
    slide_11_closing(prs)

    # Fix all page numbers
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        t = run.text.strip()
                        if '/' in t and t.replace(' ', '').replace('/', '').isdigit():
                            run.text = f"{i+1} / {total}"

    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")
    print(f"  {len(prs.slides)} slides")
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and len(t) > 5:
                        texts.append(t)
                        break
                if texts:
                    break
        print(f"  {i+1}. {texts[0][:60] if texts else '(empty)'}")


if __name__ == "__main__":
    main()
