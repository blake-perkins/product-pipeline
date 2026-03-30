#!/usr/bin/env python3
"""Update the executive briefing PowerPoint deck.

Updates slide 2 with new content and fixes stale references
throughout the deck (VM->VC, dashboard name, tags, etc.).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import sys
import os


def find_and_replace_text(prs, replacements):
    """Find and replace text across all slides, preserving formatting."""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
                            count += 1
    return count


def update_slide_2(slide):
    """Replace slide 2 content with new 'Built for Program Velocity' content."""
    # Remove all existing shapes except the background and nav bar
    shapes_to_keep = []
    nav_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            # Keep the navigation bar (Challenge, Solution, etc.) and page number
            if "Challenge" in text and "Solution" in text:
                nav_shapes.append(shape)
                continue
            if text.strip() in ("2 / 10",):
                nav_shapes.append(shape)
                continue
        # Remove everything else by clearing its text
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""

    # Now rebuild content in the first text shape we find
    # Find the main content shape (largest text frame)
    main_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape not in nav_shapes:
            if main_shape is None or shape.width > main_shape.width:
                main_shape = shape

    if not main_shape:
        print("  WARNING: Could not find main content shape on slide 2")
        return

    tf = main_shape.text_frame
    tf.clear()

    # Title
    p = tf.paragraphs[0]
    p.text = "Built for Program Velocity"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT

    # Blank line
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(4)

    # LEFT SIDE header
    p = tf.add_paragraph()
    p.text = "The Structural Gaps"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Gap 1
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(2)

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "No Shared Language"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x99, 0xD6)

    p = tf.add_paragraph()
    p.text = "SEs work in the model, developers work in code. Without a bridge, intent gets lost in translation at every handoff."
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.space_after = Pt(6)

    # Gap 2
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "No Reliable Handoff"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x99, 0xD6)

    p = tf.add_paragraph()
    p.text = "Artifacts shared over email or SharePoint have no version control, no change awareness, and no enforcement."
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.space_after = Pt(6)

    # Gap 3
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "No Connected Toolchain"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x99, 0xD6)

    p = tf.add_paragraph()
    p.text = "Requirements, tests, and deployment are maintained in separate tools with no automated thread between them."
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.space_after = Pt(12)

    # RIGHT SIDE header
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Why the Old Approach Falls Short"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.space_after = Pt(6)

    bullets = [
        "Requirements are living targets \u2014 this program moves faster than any manual process can track.",
        "SE/SW misalignment surfaces at integration \u2014 the most expensive point to find it.",
        "Audit evidence reconstructed from scratch each review cycle \u2014 unsustainable at this pace.",
        "Requirement changes silently invalidate tests \u2014 no alert, no flag, no enforcement.",
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = "\u2022  " + b
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p.space_after = Pt(4)

    # Bold closing line
    p = tf.add_paragraph()
    p.text = ""
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "This program demands a better approach. That\u2019s what this pipeline delivers."
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def main():
    input_path = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations", "MBSE_Pipeline_Executive_Brief.pptx")
    output_path = input_path  # overwrite

    print(f"Reading: {input_path}")
    prs = Presentation(input_path)

    # --- Global find/replace for stale references ---
    replacements = [
        # VM -> VC rename
        ("@VM:", "@VC:"),
        ("@VM-", "@VC-"),
        ("VM-01", "VC-01"),
        ("VM-02", "VC-02"),
        ("verification method", "verification criteria"),
        ("Verification Method", "Verification Criteria"),
        ("verification methods", "verification criteria"),
        ("Verification Methods", "Verification Criteria"),
        ("every VM ", "every VC "),
        ("specific VM", "specific VC"),
        ("one VM", "one VC"),
        ("uncovered VMs", "uncovered VCs"),
        ("orphaned @VM:", "orphaned @VC:"),
        ("@REQ: and @VM:", "@REQ: and @VC:"),
        # Tags
        ("@VM:SYS-REQ-001-VM-01", "@VC:SYS-REQ-001-VC-01"),
        ("@VM:SYS-REQ-001-VM-02", "@VC:SYS-REQ-001-VC-02"),
        ("SYS-REQ-001-VM-01", "SYS-REQ-001-VC-01"),
        ("SYS-REQ-001-VM-02", "SYS-REQ-001-VC-02"),
    ]

    count = find_and_replace_text(prs, replacements)
    print(f"  Replaced {count} text occurrences")

    # --- Update slide 2 ---
    print("  Updating slide 2: Built for Program Velocity")
    update_slide_2(prs.slides[1])

    # --- Save ---
    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
