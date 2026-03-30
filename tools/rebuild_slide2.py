#!/usr/bin/env python3
"""Rebuild slide 2 of the executive briefing deck from scratch."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os
import copy

INPUT = os.path.join(os.path.dirname(__file__), "..", "docs", "presentations", "MBSE_Pipeline_Executive_Brief.pptx")


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return it."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox


def set_shape_fill(shape, r, g, b):
    """Set solid fill on a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_rounded_rect(slide, left, top, width, height, r, g, b, corner_radius=100000):
    """Add a rounded rectangle with fill."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()  # no border
    # Set corner radius
    sp = shape._element
    prstGeom = sp.find(qn('a:prstGeom'), sp.nsmap) if hasattr(sp, 'nsmap') else None
    return shape


def add_text(tf, text, size=11, bold=False, color=(0xCC, 0xCC, 0xCC), space_after=4, alignment=PP_ALIGN.LEFT):
    """Add a paragraph to a text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    p.space_after = Pt(space_after)
    p.alignment = alignment
    return p


def main():
    prs = Presentation(INPUT)
    slide_layout = prs.slides[1].slide_layout  # keep same layout

    # Get the XML element for slide 2 to find its position
    slide_index = 1  # 0-indexed

    # Delete old slide 2
    rId = prs.slides._sldIdLst[slide_index].get(qn('r:id'))
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]

    # Add new slide at position 2 (insert after slide 1)
    new_slide = prs.slides.add_slide(slide_layout)

    # Move new slide to position 2 (it was added at the end)
    # Get the sldId element for the new slide
    new_sldId = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(new_sldId)
    prs.slides._sldIdLst.insert(slide_index, new_sldId)

    slide = new_slide

    # Slide dimensions
    slide_w = prs.slide_width   # 13.333 inches typically
    slide_h = prs.slide_height  # 7.5 inches

    # Background - dark navy (match existing slides)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    # ============================================
    # TITLE
    # ============================================
    title_box = add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11), Inches(0.6))
    add_text(title_box.text_frame, "Built for Program Velocity", size=32, bold=True,
             color=(0xFF, 0xFF, 0xFF), space_after=0)

    # ============================================
    # LEFT COLUMN - The Structural Gaps
    # ============================================
    left_x = Inches(0.8)
    left_w = Inches(5.5)

    # Section header
    header_box = add_textbox(slide, left_x, Inches(1.5), left_w, Inches(0.4))
    add_text(header_box.text_frame, "The Structural Gaps", size=16, bold=True,
             color=(0xFF, 0xFF, 0xFF), space_after=0)

    # Gap items
    gaps = [
        ("No Shared Language",
         "SEs work in the model, developers work in code. Without a bridge, intent gets lost in translation at every handoff."),
        ("No Reliable Handoff",
         "Artifacts shared over email or SharePoint have no version control, no change awareness, and no enforcement."),
        ("No Connected Toolchain",
         "Requirements, tests, and deployment are maintained in separate tools with no automated thread between them."),
    ]

    y_pos = Inches(2.1)
    for gap_title, gap_desc in gaps:
        gap_box = add_textbox(slide, left_x, y_pos, left_w, Inches(0.8))
        tf = gap_box.text_frame

        # Title
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = gap_title
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x99, 0xD6)
        p.space_after = Pt(2)

        # Description
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = gap_desc
        run2.font.size = Pt(10)
        run2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        p2.space_after = Pt(0)

        y_pos += Inches(0.85)

    # ============================================
    # RIGHT COLUMN - Why the Old Approach Falls Short
    # ============================================
    right_x = Inches(7.0)
    right_w = Inches(5.5)

    # Background box
    from pptx.enum.shapes import MSO_SHAPE
    box_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        right_x - Inches(0.2), Inches(1.4),
        right_w + Inches(0.4), Inches(4.2)
    )
    box_shape.fill.solid()
    box_shape.fill.fore_color.rgb = RGBColor(0x14, 0x25, 0x38)
    box_shape.line.fill.background()

    # Section header
    rh_box = add_textbox(slide, right_x, Inches(1.6), right_w, Inches(0.4))
    add_text(rh_box.text_frame, "Why the Old Approach Falls Short", size=14, bold=True,
             color=(0xFF, 0xFF, 0xFF), space_after=0)

    # Bullet points
    bullets = [
        "Requirements are living targets \u2014 this program moves faster than any manual process can track.",
        "SE/SW misalignment surfaces at integration \u2014 the most expensive point to find it.",
        "Audit evidence reconstructed from scratch each review cycle \u2014 unsustainable at this pace.",
        "Requirement changes silently invalidate tests \u2014 no alert, no flag, no enforcement.",
    ]

    bullet_box = add_textbox(slide, right_x, Inches(2.2), right_w, Inches(3.0))
    tf = bullet_box.text_frame
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = "\u2022  " + b
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        p.space_after = Pt(8)

    # ============================================
    # CLOSING LINE
    # ============================================
    close_box = add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5))
    add_text(close_box.text_frame,
             "This program demands a better approach. That\u2019s what this pipeline delivers.",
             size=14, bold=True, color=(0xFF, 0xFF, 0xFF), space_after=0,
             alignment=PP_ALIGN.LEFT)

    # ============================================
    # PAGE NUMBER
    # ============================================
    pg_box = add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3))
    add_text(pg_box.text_frame, "2 / 10", size=9, color=(0x88, 0x88, 0x88),
             space_after=0, alignment=PP_ALIGN.RIGHT)

    # Save
    prs.save(INPUT)
    print("Slide 2 rebuilt and saved.")


if __name__ == "__main__":
    main()
